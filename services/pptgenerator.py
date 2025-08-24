import os
from typing import Optional, List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER

# Helper: get a usable text_frame for the slide body
def _get_body_text_frame(slide):
    # Prefer BODY; some templates expose OBJECT instead of BODY
    for ph in slide.placeholders:
        try:
            ph_type = ph.placeholder_format.type
        except Exception:
            continue
        if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            return ph.text_frame

    # Fallback: create our own textbox
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(9.0), Inches(5.0))
    return box.text_frame

def generate_ppt(
    title: str,
    summary: str,
    bullets: List[str],
    kpis: Dict[str, str],
    filename: str,
    output_dir: str = "./outputs",
    template_path: Optional[str] = None,
) -> str:
    os.makedirs(output_dir, exist_ok=True)
    path = os.path.join(output_dir, filename)

    # Load template if provided; else use default theme
    prs = Presentation(template_path) if template_path else Presentation()

    # -------- Title slide (layout 0 usually works, but be defensive)
    title_layout_idx = 0 if len(prs.slide_layouts) > 0 else 0
    slide = prs.slides.add_slide(prs.slide_layouts[title_layout_idx])

    # Title (safe)
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    else:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(9.0), Inches(1.0))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)

    # Subtitle/summary — try a SUBTITLE/BODY placeholder, else add textbox
    body_tf = None
    for ph in slide.placeholders:
        try:
            ph_type = ph.placeholder_format.type
        except Exception:
            continue
        if ph_type in (PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            body_tf = ph.text_frame
            break
    if body_tf is None:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(9.0), Inches(1.6))
        body_tf = tb.text_frame
    body_tf.clear()
    body_tf.paragraphs[0].text = summary or "Summary"

    # -------- KPI slide
    content_layout_idx = 1 if len(prs.slide_layouts) > 1 else 0  # "Title and Content" if available
    slide = prs.slides.add_slide(prs.slide_layouts[content_layout_idx])

    if slide.shapes.title is not None:
        slide.shapes.title.text = "Key Metrics (Last 7 days)"
    else:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(9.0), Inches(1.0))
        p = tb.text_frame.paragraphs[0]
        p.text = "Key Metrics (Last 7 days)"
        p.font.size = Pt(28)

    tf = _get_body_text_frame(slide)
    tf.clear()
    # Write KPIs
    if kpis:
        first = True
        for k, v in kpis.items():
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            para.text = f"{k}: {v}"
            para.level = 0
    else:
        tf.paragraphs[0].text = "No KPIs available."

    # -------- Bottlenecks slide
    slide = prs.slides.add_slide(prs.slide_layouts[content_layout_idx])
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Bottlenecks & Actions"
    else:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(9.0), Inches(1.0))
        p = tb.text_frame.paragraphs[0]
        p.text = "Bottlenecks & Actions"
        p.font.size = Pt(28)

    tf = _get_body_text_frame(slide)
    tf.clear()
    if bullets:
        first = True
        for b in bullets:
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            para.text = str(b)
            para.level = 0
    else:
        tf.paragraphs[0].text = "No bottlenecks found."

    # -------- Next Steps slide (use a content layout if available; index 5 may not exist in custom templates)
    next_steps_layout_idx = content_layout_idx if len(prs.slide_layouts) <= 5 else 5
    slide = prs.slides.add_slide(prs.slide_layouts[next_steps_layout_idx])
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Next Steps"
    else:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(9.0), Inches(1.0))
        p = tb.text_frame.paragraphs[0]
        p.text = "Next Steps"
        p.font.size = Pt(28)

    body = _get_body_text_frame(slide)
    body.clear()
    default_actions = [
        "1) Validate data assumptions with process owners",
        "2) Pilot quick wins (SLA alerts, approval auto-routing)",
        "3) Re-measure cycle time & rework after 2 weeks",
    ]
    first = True
    for line in default_actions:
        if first:
            para = body.paragraphs[0]
            first = False
        else:
            para = body.add_paragraph()
        para.text = line
        para.level = 0

    prs.save(path)
    return path
